"""Generate user-requested Word, Excel, and PowerPoint files.

This service is intentionally database-light: generated files are written to a
tenant-scoped directory with sidecar metadata. That keeps the MVP simple while
preserving tenant checks at download time and leaves a clear migration path to
object storage later.
"""

from __future__ import annotations

import json
import re
import uuid
import asyncio
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Literal

from pydantic import BaseModel, Field

from app.config import settings
from app.llm.client import (
    ModelTier,
    begin_usage_collection,
    complete_structured,
    end_usage_collection,
)

GeneratedFileFormat = Literal["docx", "xlsx", "pptx"]

_FORMAT_LABELS: dict[GeneratedFileFormat, str] = {
    "docx": "Word 文档",
    "xlsx": "Excel 表格",
    "pptx": "PPT 演示文稿",
}
_MIME_TYPES: dict[GeneratedFileFormat, str] = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
_REQUEST_VERBS = ("生成", "制作", "创建", "导出", "输出", "整理", "写一份", "做一份", "帮我")
_PLAN_TIMEOUT_SECONDS = 30.0


class FileGenerationError(RuntimeError):
    """Raised when a generated file cannot be created or read."""


class FileSection(BaseModel):
    heading: str = Field(default="未命名章节")
    summary: str = ""
    bullets: list[str] = Field(default_factory=list)


class FileTable(BaseModel):
    title: str = Field(default="关键数据")
    headers: list[str] = Field(default_factory=lambda: ["维度", "内容"])
    rows: list[list[str]] = Field(default_factory=list)


class FileSlide(BaseModel):
    title: str = Field(default="未命名页面")
    bullets: list[str] = Field(default_factory=list)
    speaker_note: str = ""


class FilePlan(BaseModel):
    title: str = Field(default="AtomCAP 投研材料")
    subtitle: str = ""
    sections: list[FileSection] = Field(default_factory=list)
    tables: list[FileTable] = Field(default_factory=list)
    slides: list[FileSlide] = Field(default_factory=list)


@dataclass(frozen=True)
class GeneratedFile:
    file_id: str
    filename: str
    title: str
    format: GeneratedFileFormat
    mime_type: str
    size_bytes: int
    download_url: str
    created_at: str
    path: Path

    def to_ref(self) -> dict[str, object]:
        return {
            "type": "generated_file",
            "file_id": self.file_id,
            "filename": self.filename,
            "title": self.title,
            "format": self.format,
            "mime_type": self.mime_type,
            "size_bytes": self.size_bytes,
            "download_url": self.download_url,
            "created_at": self.created_at,
        }


@dataclass(frozen=True)
class StoredGeneratedFile:
    path: Path
    filename: str
    mime_type: str
    metadata: dict[str, object]


@dataclass(frozen=True)
class FileGenerationResult:
    file: GeneratedFile
    usage: dict[str, int] | None
    plan: FilePlan


def detect_file_generation_request(content: str) -> GeneratedFileFormat | None:
    """Return the requested output format when the prompt clearly asks for a file."""
    text = (content or "").strip()
    if not text:
        return None
    lowered = text.lower()
    asks_for_file = any(verb in text for verb in _REQUEST_VERBS)
    if not asks_for_file:
        return None
    if any(token in lowered for token in ("ppt", "pptx", "powerpoint")) or any(
        token in text for token in ("路演PPT", "幻灯片", "演示文稿", "路演材料")
    ):
        return "pptx"
    if any(token in lowered for token in ("excel", "xlsx", "xls")) or any(
        token in text for token in ("电子表格", "表格文件")
    ):
        return "xlsx"
    if any(token in lowered for token in ("word", "docx", "doc")) or "Word" in text:
        return "docx"
    return None


async def generate_file_from_request(
    *,
    institution_id: uuid.UUID,
    user_request: str,
    target_format: GeneratedFileFormat,
    runtime_context: str,
    tier: ModelTier = ModelTier.PREMIUM,
    allow_overseas: bool = False,
) -> FileGenerationResult:
    """Create a generated file and return a frontend-safe reference."""
    usage_token, usage_events = begin_usage_collection()
    try:
        plan = await _build_file_plan(
            user_request=user_request,
            target_format=target_format,
            runtime_context=runtime_context,
            tier=tier,
            allow_overseas=allow_overseas,
        )
    finally:
        end_usage_collection(usage_token)

    if not plan.sections:
        plan.sections = _fallback_plan(user_request, target_format, runtime_context).sections
    if target_format == "pptx" and not plan.slides:
        plan.slides = _slides_from_sections(plan)
    if target_format == "xlsx" and not plan.tables:
        plan.tables = _tables_from_sections(plan)

    generated = create_generated_file_from_plan(
        institution_id=institution_id,
        plan=plan,
        target_format=target_format,
    )
    return FileGenerationResult(
        file=generated,
        usage=_aggregate_usage(usage_events),
        plan=plan,
    )


def create_generated_file_from_plan(
    *,
    institution_id: uuid.UUID,
    plan: FilePlan,
    target_format: GeneratedFileFormat,
) -> GeneratedFile:
    """Render a pre-built file plan without calling the LLM."""
    root = _tenant_file_dir(institution_id)
    root.mkdir(parents=True, exist_ok=True)
    file_id = str(uuid.uuid4())
    filename = f"{_safe_filename(plan.title)}.{target_format}"
    file_path = root / f"{file_id}.{target_format}"

    if target_format == "docx":
        _render_docx(plan, file_path)
    elif target_format == "xlsx":
        _render_xlsx(plan, file_path)
    elif target_format == "pptx":
        _render_pptx(plan, file_path)
    else:  # pragma: no cover - Literal keeps this unreachable.
        raise FileGenerationError(f"不支持的文件格式：{target_format}")

    created_at = datetime.now(timezone.utc).isoformat()
    generated = GeneratedFile(
        file_id=file_id,
        filename=filename,
        title=plan.title,
        format=target_format,
        mime_type=_MIME_TYPES[target_format],
        size_bytes=file_path.stat().st_size,
        download_url=f"/api/conversations/files/{file_id}",
        created_at=created_at,
        path=file_path,
    )
    _write_metadata(institution_id=institution_id, generated=generated)
    return generated


def get_generated_file(
    *, institution_id: uuid.UUID, file_id: uuid.UUID | str
) -> StoredGeneratedFile:
    """Load metadata and file path for a tenant-owned generated file."""
    file_uuid = str(uuid.UUID(str(file_id)))
    root = _tenant_file_dir(institution_id)
    metadata_path = root / f"{file_uuid}.json"
    if not metadata_path.exists():
        raise FileGenerationError("文件不存在或已被清理。")
    metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
    if str(metadata.get("institution_id")) != str(institution_id):
        raise FileGenerationError("无权访问该文件。")
    target_format = metadata.get("format")
    if target_format not in _MIME_TYPES:
        raise FileGenerationError("文件元数据格式无效。")
    file_path = root / f"{file_uuid}.{target_format}"
    if not file_path.exists():
        raise FileGenerationError("文件实体不存在或已被清理。")
    return StoredGeneratedFile(
        path=file_path,
        filename=str(metadata.get("filename") or file_path.name),
        mime_type=str(metadata.get("mime_type") or _MIME_TYPES[target_format]),
        metadata=metadata,
    )


async def _build_file_plan(
    *,
    user_request: str,
    target_format: GeneratedFileFormat,
    runtime_context: str,
    tier: ModelTier,
    allow_overseas: bool,
) -> FilePlan:
    system = (
        "你是 AtomCAP 的文件生成工具，负责把投研对话和工作台上下文整理成可落地的文件内容。"
        "你只需要输出结构化内容，不要输出 Markdown，不要虚构不存在的事实。"
        "如果上下文不足，必须在内容中标明“待补充/待核验”。"
    )
    format_hint = {
        "docx": "生成适合 Word 投研报告的章节，建议 5-8 个章节，每章 3-6 个要点。",
        "xlsx": "生成适合 Excel 的结构化表格，至少包含摘要表、关键事实表和待办/风险表。",
        "pptx": "生成适合路演或内部汇报 PPT 的页面，建议 8-12 页，每页 3-5 个要点。",
    }[target_format]
    user = "\n\n".join(
        part
        for part in (
            f"用户请求：{user_request}",
            f"目标格式：{_FORMAT_LABELS[target_format]}",
            f"生成要求：{format_hint}",
            "可用上下文：\n" + (runtime_context.strip() or "暂无额外上下文。"),
        )
        if part.strip()
    )
    try:
        return await asyncio.wait_for(
            complete_structured(
                tier,
                [{"role": "system", "content": system}, {"role": "user", "content": user}],
                FilePlan,
                allow_overseas=allow_overseas,
                max_repair_attempts=1,
            ),
            timeout=_PLAN_TIMEOUT_SECONDS,
        )
    except Exception as exc:  # noqa: BLE001
        fallback = _fallback_plan(user_request, target_format, runtime_context)
        fallback.subtitle = f"模型生成暂不可用，已使用确定性模板兜底：{type(exc).__name__}"
        return fallback


def _fallback_plan(
    user_request: str, target_format: GeneratedFileFormat, runtime_context: str
) -> FilePlan:
    title = _infer_title(user_request, target_format)
    context_excerpt = _compact(runtime_context, limit=500) or "当前没有额外上下文。"
    sections = [
        FileSection(
            heading="任务摘要",
            summary=user_request.strip() or "用户要求生成投研文件。",
            bullets=[
                "本文件由 AtomCAP 文件生成工具根据当前会话和工作台上下文生成。",
                "未能从上下文确认的事实均应在后续投研中继续核验。",
            ],
        ),
        FileSection(
            heading="已读取上下文",
            summary="系统可用上下文摘要。",
            bullets=[context_excerpt],
        ),
        FileSection(
            heading="建议补充",
            summary="为提高文件质量，建议继续补充以下材料。",
            bullets=["项目或赛道的最新公开资料", "机构内部判断和关键假设", "可验证的数据来源和证据链接"],
        ),
    ]
    return FilePlan(
        title=title,
        subtitle="AtomCAP 自动生成",
        sections=sections,
        tables=_tables_from_sections(FilePlan(title=title, sections=sections)),
        slides=_slides_from_sections(FilePlan(title=title, sections=sections)),
    )


def _render_docx(plan: FilePlan, path: Path) -> None:
    try:
        from docx import Document as WordDocument
    except Exception as exc:  # noqa: BLE001
        raise FileGenerationError("生成 Word 需要安装 python-docx。") from exc

    doc = WordDocument()
    doc.add_heading(plan.title or "AtomCAP 投研材料", 0)
    if plan.subtitle:
        doc.add_paragraph(plan.subtitle)
    doc.add_paragraph(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}")
    for section in plan.sections:
        doc.add_heading(section.heading or "未命名章节", level=1)
        if section.summary:
            doc.add_paragraph(section.summary)
        for bullet in section.bullets:
            doc.add_paragraph(bullet, style="List Bullet")
    for table in plan.tables:
        doc.add_heading(table.title or "数据表", level=1)
        headers = table.headers or ["维度", "内容"]
        rows = table.rows or []
        word_table = doc.add_table(rows=max(len(rows) + 1, 1), cols=max(len(headers), 1))
        word_table.style = "Table Grid"
        for idx, header in enumerate(headers):
            word_table.cell(0, idx).text = str(header)
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, value in enumerate(row[: len(headers)]):
                word_table.cell(row_idx, col_idx).text = str(value)
    doc.save(path)


def _render_xlsx(plan: FilePlan, path: Path) -> None:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
    except Exception as exc:  # noqa: BLE001
        raise FileGenerationError("生成 Excel 需要安装 openpyxl。") from exc

    wb = Workbook()
    summary_ws = wb.active
    summary_ws.title = "摘要"
    title_fill = PatternFill("solid", fgColor="EEF2FF")
    summary_ws["A1"] = plan.title or "AtomCAP 投研材料"
    summary_ws["A1"].font = Font(size=16, bold=True)
    summary_ws["A2"] = plan.subtitle
    row = 4
    for section in plan.sections:
        summary_ws.cell(row=row, column=1, value=section.heading).font = Font(bold=True)
        summary_ws.cell(row=row, column=1).fill = title_fill
        row += 1
        if section.summary:
            summary_ws.cell(row=row, column=1, value=section.summary)
            row += 1
        for bullet in section.bullets:
            summary_ws.cell(row=row, column=1, value=bullet)
            row += 1
        row += 1
    summary_ws.column_dimensions["A"].width = 80
    summary_ws["A1"].alignment = Alignment(wrap_text=True)

    for idx, table in enumerate(plan.tables or _tables_from_sections(plan), start=1):
        ws = wb.create_sheet(_safe_sheet_name(table.title or f"表{idx}"))
        headers = table.headers or ["维度", "内容"]
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = title_fill
        for record in table.rows:
            ws.append([str(value) for value in record[: len(headers)]])
        for column_cells in ws.columns:
            column_letter = column_cells[0].column_letter
            ws.column_dimensions[column_letter].width = min(
                max(max(len(str(cell.value or "")) for cell in column_cells) + 4, 14),
                48,
            )
            for cell in column_cells:
                cell.alignment = Alignment(wrap_text=True, vertical="top")
    wb.save(path)


def _render_pptx(plan: FilePlan, path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as exc:  # noqa: BLE001
        raise FileGenerationError("生成 PPT 需要安装 python-pptx，请先刷新后端依赖。") from exc

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = plan.title or "AtomCAP 投研材料"
    title_slide.placeholders[1].text = plan.subtitle or "自动生成"

    for slide in plan.slides or _slides_from_sections(plan):
        layout = prs.slide_layouts[1]
        page = prs.slides.add_slide(layout)
        page.shapes.title.text = slide.title or "未命名页面"
        body = page.shapes.placeholders[1].text_frame
        body.clear()
        bullets = slide.bullets or ["待补充"]
        for idx, bullet in enumerate(bullets[:6]):
            para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            para.text = str(bullet)
            para.level = 0
            para.font.size = Pt(20)
        if slide.speaker_note:
            notes = page.notes_slide.notes_text_frame
            notes.text = slide.speaker_note
    prs.save(path)


def _slides_from_sections(plan: FilePlan) -> list[FileSlide]:
    slides = [
        FileSlide(title=section.heading, bullets=section.bullets or [section.summary or "待补充"])
        for section in plan.sections
    ]
    return slides or [FileSlide(title=plan.title, bullets=["待补充上下文后完善内容。"])]


def _tables_from_sections(plan: FilePlan) -> list[FileTable]:
    rows: list[list[str]] = []
    for section in plan.sections:
        content = section.summary or "；".join(section.bullets[:3]) or "待补充"
        rows.append([section.heading, content])
    return [FileTable(title="内容摘要", headers=["模块", "摘要"], rows=rows)]


def _write_metadata(*, institution_id: uuid.UUID, generated: GeneratedFile) -> None:
    metadata = {
        **generated.to_ref(),
        "institution_id": str(institution_id),
        "stored_path": str(generated.path),
    }
    metadata_path = generated.path.with_suffix(".json")
    metadata_path.write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")


def _tenant_file_dir(institution_id: uuid.UUID) -> Path:
    base = Path(settings.generated_files_dir).expanduser()
    if not base.is_absolute():
        base = Path(__file__).resolve().parents[2] / base
    return base / str(institution_id)


def _aggregate_usage(events: list[dict[str, int]]) -> dict[str, int] | None:
    if not events:
        return None
    total: dict[str, int] = {}
    for event in events:
        for key, value in event.items():
            total[key] = total.get(key, 0) + int(value)
    return total


def _safe_filename(value: str, *, fallback: str = "AtomCAP文件") -> str:
    name = re.sub(r"[\\/:*?\"<>|\r\n\t]+", " ", value or "").strip()
    name = re.sub(r"\s+", " ", name)
    return (name or fallback)[:80]


def _safe_sheet_name(value: str) -> str:
    name = re.sub(r"[\[\]:*?/\\]+", " ", value or "").strip() or "Sheet"
    return name[:31]


def _infer_title(user_request: str, target_format: GeneratedFileFormat) -> str:
    compact = _compact(user_request, limit=36) or "投研材料"
    suffix = {
        "docx": "Word 文档",
        "xlsx": "Excel 表格",
        "pptx": "路演 PPT",
    }[target_format]
    return f"{compact} - {suffix}"


def _compact(value: str, *, limit: int = 800) -> str:
    text = " ".join((value or "").split())
    return text if len(text) <= limit else f"{text[:limit]}..."
